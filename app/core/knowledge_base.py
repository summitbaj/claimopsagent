"""
Knowledge Base management using ChromaDB for storing and retrieving documents.
Supports ingestion of PowerPoint (PPTX), PDF, and Word (DOCX) files.
"""
import os
import json
from typing import List, Dict, Any
from pathlib import Path
import chromadb
from chromadb.config import Settings
from pptx import Presentation
from pypdf import PdfReader
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from app.core.config import settings
import logging
from datetime import datetime

logger = logging.getLogger(__name__)


class KnowledgeBase:
    """
    Manages document ingestion and retrieval using ChromaDB.
    """
    
    def __init__(self, persist_directory: str = "data/chroma_db"):
        """
        Initialize ChromaDB client.
        
        Args:
            persist_directory: Directory to store ChromaDB data
        """
        self.persist_directory = persist_directory
        
        # Create directory if it doesn't exist
        Path(persist_directory).mkdir(parents=True, exist_ok=True)
        
        # Initialize ChromaDB client (v0.3.23 compatible)
        self.client = chromadb.Client(Settings(
            chroma_db_impl="duckdb+parquet",
            persist_directory=persist_directory,
            anonymized_telemetry=False
        ))

        # Use default embedding function explicitly
        from chromadb.utils import embedding_functions
        self.embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction()
        
        # Get or create collection
        self.collection = self.client.get_or_create_collection(
            name="knowledge_base",
            embedding_function=self.embedding_fn,
            metadata={"description": "Claims and billing operations knowledge"}
        )
        
        logger.info(f"📚 KnowledgeBase initialized at {persist_directory}")
    
    def _extract_text_from_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from PowerPoint file.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            List of dictionaries with slide text and metadata
        """
        presentation = Presentation(file_path)
        slides_data = []
        
        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            if slide_text:
                slides_data.append({
                    "text": "\n".join(slide_text),
                    "slide_number": slide_num,
                    "source": os.path.basename(file_path)
                })
        
        logger.info(f"📄 Extracted {len(slides_data)} slides from {file_path}")
        return slides_data
    
    def _extract_text_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from PDF file.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            List of dictionaries with page text and metadata
        """
        reader = PdfReader(file_path)
        pages_data = []
        
        for page_num, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                pages_data.append({
                    "text": text,
                    "page_number": page_num,
                    "source": os.path.basename(file_path)
                })
        
        logger.info(f"📄 Extracted {len(pages_data)} pages from {file_path}")
        return pages_data
    
    def _extract_text_from_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from Word file.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            List of dictionaries with paragraph text and metadata
        """
        doc = Document(file_path)
        paragraphs_data = []
        
        # Group paragraphs into logical sections
        current_section = []
        section_num = 1
        
        for para in doc.paragraphs:
            if para.text.strip():
                current_section.append(para.text)
                # Create a section every 5 paragraphs or when encountering a heading
                if len(current_section) >= 5 or (para.style.name.startswith('Heading')):
                    if current_section:
                        paragraphs_data.append({
                            "text": "\n".join(current_section),
                            "section_number": section_num,
                            "source": os.path.basename(file_path)
                        })
                        section_num += 1
                        current_section = []
        
        # Add remaining paragraphs
        if current_section:
            paragraphs_data.append({
                "text": "\n".join(current_section),
                "section_number": section_num,
                "source": os.path.basename(file_path)
            })
        
        logger.info(f"📄 Extracted {len(paragraphs_data)} sections from {file_path}")
        return paragraphs_data
    
    def ingest_file(self, file_path: str, file_type: str = None) -> Dict[str, Any]:
        """
        Ingest a file (PPTX, PDF, or DOCX) into the knowledge base.
        
        Args:
            file_path: Path to the file
            file_type: Optional file type override (pptx, pdf, docx)
            
        Returns:
            Dictionary with ingestion results
        """
        try:
            # Determine file type
            if file_type is None:
                file_type = Path(file_path).suffix.lower().lstrip('.')
            
            # Extract text based on file type
            if file_type in ['pptx', 'ppt']:
                extracted_data = self._extract_text_from_pptx(file_path)
                unit_name = "slides"
                unit_key = "slide_number"
            elif file_type == 'pdf':
                extracted_data = self._extract_text_from_pdf(file_path)
                unit_name = "pages"
                unit_key = "page_number"
            elif file_type in ['docx', 'doc']:
                extracted_data = self._extract_text_from_docx(file_path)
                unit_name = "sections"
                unit_key = "section_number"
            else:
                return {
                    "success": False,
                    "error": f"Unsupported file type: {file_type}"
                }
            
            if not extracted_data:
                return {
                    "success": False,
                    "error": "No text content found in file"
                }
            
            # Split text into chunks
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len
            )
            
            documents = []
            metadatas = []
            ids = []
            
            filename = os.path.basename(file_path)
            doc_id_base = f"{filename}_{self.collection.count()}"
            
            for data in extracted_data:
                chunks = text_splitter.split_text(data["text"])
                
                for chunk_idx, chunk in enumerate(chunks):
                    doc_id = f"{doc_id_base}_{unit_name[:4]}{data[unit_key]}_chunk{chunk_idx}"
                    
                    documents.append(chunk)
                    metadatas.append({
                        "source": data["source"],
                        "file_type": file_type,
                        unit_key: str(data[unit_key]),
                        "chunk_index": str(chunk_idx)
                    })
                    ids.append(doc_id)
            
            # Add to ChromaDB
            self.collection.add(
                documents=documents,
                metadatas=metadatas,
                ids=ids
            )
            
            # Persist changes
            self.client.persist()
            
            # Save source metadata
            self._save_source_metadata(filename, file_type, len(extracted_data), len(documents))
            
            logger.info(f"✅ Ingested {len(documents)} chunks from {file_path}")
            
            return {
                "success": True,
                "file": filename,
                "file_type": file_type,
                "units_processed": len(extracted_data),
                "chunks_created": len(documents)
            }
            
        except Exception as e:
            logger.error(f"❌ Error ingesting file: {str(e)}", exc_info=True)
            return {
                "success": False,
                "error": str(e)
            }
    
    def _save_source_metadata(self, filename: str, file_type: str, units: int, chunks: int):
        """Save metadata about uploaded sources."""
        metadata_dir = Path("data/metadata")
        metadata_dir.mkdir(parents=True, exist_ok=True)
        metadata_file = metadata_dir / "sources.json"
        
        # Load existing metadata
        sources = []
        if metadata_file.exists():
            with open(metadata_file, 'r') as f:
                sources = json.load(f)
        
        # Add new source
        sources.append({
            "filename": filename,
            "file_type": file_type,
            "units_processed": units,
            "chunks_created": chunks,
            "uploaded_at": datetime.now().isoformat()
        })
        
        # Save metadata
        with open(metadata_file, 'w') as f:
            json.dump(sources, f, indent=2)
    
    def get_uploaded_sources(self) -> List[Dict[str, Any]]:
        """Get list of all uploaded knowledge sources."""
        metadata_file = Path("data/metadata/sources.json")
        if metadata_file.exists():
            with open(metadata_file, 'r') as f:
                return json.load(f)
        return []
    
    def delete_source(self, filename: str) -> Dict[str, Any]:
        """
        Delete a knowledge source and its chunks from the database.
        
        Args:
            filename: Name of the file to delete
            
        Returns:
            Dictionary with deletion results
        """
        try:
            # Get all documents with this source
            results = self.collection.get(
                where={"source": filename}
            )
            
            if results and results['ids']:
                # Delete the documents
                self.collection.delete(ids=results['ids'])
                
                # Persist changes
                self.client.persist()
                
                deleted_count = len(results['ids'])
                
                # Update metadata
                metadata_file = Path("data/metadata/sources.json")
                if metadata_file.exists():
                    with open(metadata_file, 'r') as f:
                        sources = json.load(f)
                    
                    sources = [s for s in sources if s['filename'] != filename]
                    
                    with open(metadata_file, 'w') as f:
                        json.dump(sources, f, indent=2)
                
                logger.info(f"🗑️ Deleted {deleted_count} chunks for {filename}")
                return {
                    "success": True,
                    "chunks_deleted": deleted_count
                }
            else:
                return {
                    "success": False,
                    "error": "Source not found"
                }
                
        except Exception as e:
            logger.error(f"❌ Error deleting source: {str(e)}", exc_info=True)
            return {
                "success": False,
                "error": str(e)
            }
    
    def query(self, query_text: str, n_results: int = 3) -> List[Dict[str, Any]]:
        """
        Query the knowledge base for relevant documents.
        
        Args:
            query_text: The query string
            n_results: Number of results to return
            
        Returns:
            List of relevant documents with metadata
        """
        try:
            results = self.collection.query(
                query_texts=[query_text],
                n_results=n_results
            )
            
            # Format results
            formatted_results = []
            
            if results and results['documents'] and results['documents'][0]:
                for i, doc in enumerate(results['documents'][0]):
                    metadata = results['metadatas'][0][i] if results['metadatas'] else {}
                    distance = results['distances'][0][i] if results['distances'] else None
                    
                    # Get page/slide/section number
                    location = metadata.get("slide_number") or metadata.get("page_number") or metadata.get("section_number", "N/A")
                    
                    formatted_results.append({
                        "content": doc,
                        "source": metadata.get("source", "Unknown"),
                        "slide_number": location,  # Keep for backward compatibility
                        "relevance_score": 1 - distance if distance is not None else None
                    })
            
            logger.info(f"🔍 Query returned {len(formatted_results)} results")
            return formatted_results
            
        except Exception as e:
            logger.error(f"❌ Error querying knowledge base: {str(e)}", exc_info=True)
            return []
    
    def get_collection_stats(self) -> Dict[str, Any]:
        """
        Get statistics about the knowledge base collection.
        
        Returns:
            Dictionary with collection statistics
        """
        try:
            count = self.collection.count()
            return {
                "total_documents": count,
                "collection_name": self.collection.name,
                "persist_directory": self.persist_directory
            }
        except Exception as e:
            logger.error(f"❌ Error getting collection stats: {str(e)}")
            return {"error": str(e)}
